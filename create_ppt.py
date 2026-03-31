from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = "Times New Roman"
    p.font.color.rgb = RGBColor(102, 126, 234)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.name = "Times New Roman"
        p2.font.color.rgb = RGBColor(148, 163, 184)
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
    shape.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(102, 126, 234)
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Times New Roman"
    p.font.color.rgb = RGBColor(248, 250, 252)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list, bullets=True):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
    shape.line.fill.background()
    
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(30, 41, 59)
    header_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Times New Roman"
    p.font.color.rgb = RGBColor(102, 126, 234)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        
        if bullets and not item.startswith("•"):
            p.text = "• " + item
        else:
            p.text = item
        
        p.font.size = Pt(18)
        p.font.name = "Times New Roman"
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_after = Pt(12)
        p.line_spacing = 1.5
    
    return slide

def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Times New Roman"
    p.font.color.rgb = RGBColor(102, 126, 234)
    
    num_cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, num_cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5)).table
    
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(102, 126, 234)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = "Times New Roman"
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.name = "Times New Roman"
            p.font.color.rgb = RGBColor(226, 232, 240)
    
    return slide

# SLIDE 1: Title
add_title_slide(prs, "TaskMind", "AI-Powered Task Scheduler")
add_title_slide(prs, "Project Presentation", "Advanced Python Scheduler (APScheduler) with AI Priority Scheduling")

# SLIDE 2: Table of Contents
toc_content = [
    "1. Abstract",
    "2. Introduction",
    "3. Need for the Study",
    "4. Objectives",
    "5. Literature Survey",
    "6. Methodology Used to Achieve Objectives",
    "7. Software Requirements",
    "8. Execution Screenshots",
    "9. Code in Separated Files",
    "10. Creativity and Innovation",
    "11. Future Scope",
    "12. Conclusion"
]
add_content_slide(prs, "Table of Contents", toc_content)

# SLIDE 3: Abstract
abstract = [
    "TaskMind is an AI-powered task scheduling system that automatically prioritizes and executes tasks based on multiple factors including deadline urgency, task complexity, and execution history.",
    "Unlike traditional OS schedulers that use simple First-Come-First-Serve or round-robin algorithms, TaskMind employs a sophisticated AI scoring mechanism that dynamically calculates priority scores for each task.",
    "The system provides real-time monitoring of system resources (CPU and Memory), visual dashboards for task management, and an intuitive web interface for task submission and monitoring.",
    "This project demonstrates the application of AI in optimizing task scheduling, which has significant real-world applications in healthcare, financial trading, e-commerce, and data centers.",
    "The implementation uses modern web technologies with a Python backend, providing both scalability and ease of deployment across various platforms."
]
add_content_slide(prs, "Abstract", abstract)

# SLIDE 4: Introduction
intro = [
    "Task scheduling is a fundamental concept in operating systems that determines which process runs at what time on the CPU.",
    "Traditional schedulers like FCFS, SJF, Round Robin, and Priority Scheduling have limitations in handling real-world scenarios where multiple factors influence task urgency.",
    "TaskMind addresses these limitations by implementing an AI-based priority scoring system that considers deadline proximity, task complexity, and historical execution performance.",
    "The system acts as a layer above traditional OS scheduling, providing intelligent task prioritization for applications requiring complex scheduling decisions.",
    "With the increasing complexity of modern computing workloads, AI-powered scheduling becomes essential for optimizing resource utilization and meeting service level agreements."
]
add_content_slide(prs, "Introduction", intro)

# SLIDE 5: Need for the Study
need = [
    "Modern applications require intelligent scheduling beyond simple time-based prioritization",
    "Real-time systems need to balance multiple competing priorities simultaneously",
    "Existing scheduling algorithms lack adaptability and learning capabilities",
    "Growing demand for personalized task management in enterprise environments",
    "Need for transparent, explainable scheduling decisions in critical applications",
    "Gap between theoretical OS scheduling concepts and practical AI applications",
    "Requirement for visual dashboards to monitor scheduling performance in real-time",
    "Demand for scalable solutions that can handle both simple and complex scheduling scenarios"
]
add_content_slide(prs, "Need for the Study", need)

# SLIDE 6: Objectives
objectives = [
    "To develop an AI-powered task scheduling system using APScheduler",
    "To implement dynamic priority calculation based on multiple factors",
    "To create a real-time resource monitoring dashboard for CPU and Memory",
    "To design an intuitive web interface for task management",
    "To demonstrate OS scheduling concepts through practical implementation",
    "To provide task execution history tracking and analytics",
    "To implement task dependencies and DAG-based scheduling",
    "To create a system that can be deployed on both local and cloud platforms"
]
add_content_slide(prs, "Objectives", objectives)

# SLIDE 7: Literature Survey
literature = [
    "Traditional OS schedulers (FCFS, SJF, Round Robin) - Basis for all modern scheduling",
    "APScheduler (Advanced Python Scheduler) - Industry-standard Python task scheduling library",
    "Priority Inversion and Inheritance Protocols - Handle priority conflicts in real systems",
    "Multi-level Feedback Queue Scheduling - Adaptive algorithm using historical data",
    "Earliest Deadline First (EDF) - Real-time scheduling based on time constraints",
    "Machine Learning in Job Scheduling - Research shows 40-60% improvement in resource utilization",
    "Cloud-based Task Orchestration (Kubernetes, Airflow) - Modern distributed scheduling patterns"
]
add_content_slide(prs, "Literature Survey", literature)

# SLIDE 8: Methodology Table
headers = ["Phase", "Technology", "Purpose"]
rows = [
    ["Frontend", "HTML5, CSS3, Tailwind CSS", "User Interface Design"],
    ["Charts", "Chart.js", "Real-time Resource Visualization"],
    ["Backend API", "Python FastAPI", "REST API Development"],
    ["AI Engine", "Custom Algorithm", "Priority Score Calculation"],
    ["Task Queue", "APScheduler", "Task Scheduling Engine"],
    ["Server", "Uvicorn", "ASGI Server"],
    ["Local Deploy", "Node.js", "Local Development Server"],
    ["Cloud Deploy", "Vercel", "Production Deployment"]
]
add_table_slide(prs, "Methodology - Technology Stack", headers, rows)

# SLIDE 9: Methodology - AI Priority Algorithm
methodology = [
    "AI Priority Score = (Deadline Factor × 0.5) + (Complexity Factor × 0.3) + (History Factor × 0.2)",
    "Deadline Factor: Calculated based on hours remaining until deadline (0-100 scale)",
    "Complexity Factor: User-provided complexity value scaled by 10 (1-10 input = 10-100 output)",
    "History Factor: Based on past execution success rate and average duration",
    "Tasks are sorted in descending order of priority score for execution",
    "Real-time updates recalculate priorities as conditions change",
    "Resource monitoring runs in parallel, tracking CPU and memory usage",
    "Dashboard provides visual feedback of all scheduling operations"
]
add_content_slide(prs, "Methodology - AI Priority Algorithm", methodology)

# SLIDE 10: Software Requirements
software = [
    "Python 3.8+ - Programming Language",
    "FastAPI - Modern web framework for API development",
    "APScheduler - Task scheduling library",
    "Uvicorn - ASGI server for running FastAPI",
    "Node.js - For local server deployment",
    "Tailwind CSS - Utility-first CSS framework",
    "Chart.js - JavaScript library for charts",
    "Git - Version control system",
    "Visual Studio Code / PyCharm - IDE",
    "Chrome/Firefox Browser - For testing"
]
add_content_slide(prs, "Software Requirements", software)

# SLIDE 11: Hardware Requirements
hardware = [
    "Processor: Intel Core i3 or equivalent",
    "RAM: Minimum 4GB (Recommended 8GB)",
    "Storage: 1GB available space",
    "Network: Internet connection for cloud deployment",
    "Display: 1280x720 minimum resolution",
    "OS: Windows 10+, macOS, or Linux"
]
add_content_slide(prs, "Hardware Requirements", hardware)

# SLIDE 12: Execution Screenshots
screenshots = [
    "Dashboard showing Total Tasks, Running, Completed counts",
    "Task Queue displaying all tasks with priority scores",
    "AI Priority Queue ranking tasks by calculated priority",
    "System Resources panel with CPU and Memory charts",
    "Add Task form with Name, Complexity, Duration, Deadline fields",
    "Execute Next button triggering highest priority task",
    "Run button for individual task execution",
    "Success toast notifications for completed operations"
]
add_content_slide(prs, "Execution Screenshots", screenshots)

# SLIDE 13: Code Structure
code_structure = [
    "src/ai_scheduler/__init__.py - Package initialization",
    "src/ai_scheduler/web.py - FastAPI web application and HTML dashboard",
    "src/ai_scheduler/priority_scorer.py - AI priority calculation algorithm",
    "src/ai_scheduler/resource_monitor.py - CPU/Memory monitoring",
    "src/ai_scheduler/dashboard.py - Task dashboard management",
    "src/ai_scheduler/storage.py - Task data persistence",
    "src/ai_scheduler/dag.py - Task dependency graph",
    "api/index.js - Vercel serverless function",
    "local-server.js - Node.js local server",
    "main.py - Python application entry point",
    "examples/ai_scheduler_demo.py - Demonstration script"
]
add_content_slide(prs, "Code Structure", code_structure)

# SLIDE 14: Key Code Snippet - Priority Algorithm
priority_code = [
    "def calculate_priority(self, deadline_factor, complexity, history_score):",
    "    # Weighted combination of factors",
    "    priority = (deadline_factor * 0.5) + (complexity * 0.3) + (history_score * 0.2)",
    "    return min(100, max(0, priority))",
    "",
    "def suggest_execution_order(self, task_ids):",
    "    # Sort tasks by priority in descending order",
    "    ranked = [(tid, self.get_priority_score(tid)) for tid in task_ids]",
    "    ranked.sort(key=lambda x: x[1], reverse=True)",
    "    return [task_id for task_id, _ in ranked]"
]
add_content_slide(prs, "Key Code - Priority Algorithm", priority_code)

# SLIDE 15: Creativity and Innovation
creativity = [
    "AI-Driven Prioritization: Combines multiple factors unlike traditional schedulers",
    "Real-time Visualization: Live charts for CPU and memory monitoring",
    "Adaptive Learning: System improves based on execution history",
    "User-Friendly Interface: Modern dark theme with intuitive controls",
    "Flexible Deployment: Works on local machines, Vercel, and Railway",
    "Task Dependencies: DAG-based scheduling for complex workflows",
    "Visual Priority Queue: Shows AI-ranked task list in real-time",
    "Success Rate Tracking: Monitors task completion statistics"
]
add_content_slide(prs, "Creativity and Innovation", creativity)

# SLIDE 16: Future Scope
future = [
    "Integration with Machine Learning models for better prediction",
    "Multi-user support with authentication and authorization",
    "Distributed scheduling across multiple nodes",
    "Support for more data stores (PostgreSQL, MongoDB)",
    "Mobile app development for iOS and Android",
    "Integration with Slack, Teams for notifications",
    "Advanced analytics and reporting dashboard",
    "Auto-scaling based on task load and resource availability",
    "Support for container orchestration (Docker, Kubernetes)",
    "Real-time collaboration features for team task management"
]
add_content_slide(prs, "Future Scope", future)

# SLIDE 17: Conclusion
conclusion = [
    "TaskMind successfully demonstrates AI-powered task scheduling with practical applications",
    "The system effectively combines deadline, complexity, and history factors for smart prioritization",
    "Real-time monitoring provides valuable insights into system resource utilization",
    "The web interface makes task management accessible to non-technical users",
    "Project showcases integration of OS scheduling concepts with modern web technologies",
    "Demonstrates practical implementation of APScheduler beyond basic usage",
    "Provides a foundation for building more sophisticated scheduling systems",
    "Successfully bridges the gap between theoretical concepts and real-world applications"
]
add_content_slide(prs, "Conclusion", conclusion)

# SLIDE 18: Thank You
add_title_slide(prs, "Thank You!", "Questions?")
add_title_slide(prs, "References", "APScheduler Docs, FastAPI Docs, Tailwind CSS, Chart.js")

prs.save("TaskMind_Project_Presentation.pptx")
print("Presentation created successfully: TaskMind_Project_Presentation.pptx")
